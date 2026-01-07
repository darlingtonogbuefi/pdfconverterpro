# backend\services\pdf_to_powerpoint.py

from pathlib import Path
import os
from pptx import Presentation
from backend.utils.file_utils import encode_file_to_base64
import pymupdf
from pptx.util import Inches
from paddleocr import PaddleOCR, PPStructureV3 as PPStructure
from PIL import Image
import io
import numpy as np
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
from sklearn.mixture import GaussianMixture
from sklearn.model_selection import GridSearchCV

class Converter:
    def __init__(self, default_font: str | None = None, enable_ocr: bool = False,
                 enforce_default_font: bool = False, image_retention_level: float = 1.0, lang: str = "ch"):
        self.default_font = default_font
        self.enable_ocr = enable_ocr  # Default to False for OCR Disabled
        self.enforce_default_font = enforce_default_font

        # Initialize PaddleOCR for OCR processing if enable_ocr is True
        self.ocr = PaddleOCR(lang=lang) if self.enable_ocr else None
        
        # PPStructure initialization (no `table` argument)
        self.layout_engine = (
            PPStructure(layout_score_threshold=image_retention_level)
            if self.enable_ocr else None
        )

    def convert(self, input_file_path: str, output_file_path: str) -> None:
        try:
            # Check if the input PDF exists
            if not os.path.isfile(input_file_path):
                raise FileNotFoundError(f"The file {input_file_path} does not exist.")
            if not input_file_path.endswith(".pdf"):
                raise ValueError("Input file must be a .pdf")
            if not output_file_path.endswith(".pptx"):
                raise ValueError("Output file must be a .pptx")

            pptx_output = Presentation()
            with pymupdf.open(input_file_path) as pdf_document:
                first_page = pdf_document[0]
                slide_width = Inches(first_page.rect.width / 72.0)
                slide_height = Inches(first_page.rect.height / 72.0)
                pptx_output.slide_width = slide_width
                pptx_output.slide_height = slide_height

                all_contents = [self._get_page_contents(page) for page in pdf_document]
                scanned_document = self.enable_ocr and all(not page["text_blocks"] for page in all_contents)
                self._construct_pptx(pdf_document, all_contents, pptx_output, scanned_document)

            # Ensure output directory exists
            output_dir = os.path.dirname(output_file_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # Save the PowerPoint
            pptx_output.save(output_file_path)

        except Exception as e:
            raise RuntimeError(f"Error during PDF to PPTX conversion: {e}")

    def _get_page_contents(self, pdf_page):
        return {
            "text_blocks": [b for b in pdf_page.get_textpage().extractDICT()["blocks"] if b["type"] == 0],
            "images": pdf_page.get_image_info(xrefs=True),
            "drawings": pdf_page.get_drawings(),
            "xref_smask_map": {item[0]: item[1] for item in pdf_page.get_images()},
        }

    def _construct_pptx(self, pdf_document, all_contents, pptx_output, scanned_document):
        if scanned_document:
            # OCR is only used if `scanned_document` is true
            ocr_results = [self.ocr.ocr(pdf_document[page_num].get_pixmap(dpi=300).tobytes(), cls=False)
                           for page_num in range(len(pdf_document))]
        for page_num, pdf_page in enumerate(pdf_document):
            slide = pptx_output.slides.add_slide(pptx_output.slide_layouts[6])
            page_content = all_contents[page_num]

            # Add text blocks
            for text_block in page_content["text_blocks"]:
                for line in text_block["lines"]:
                    for span in line["spans"]:
                        if span["size"] < 1:
                            continue
                        x0, y0, x1, y1 = span["bbox"]
                        textbox = slide.shapes.add_textbox(
                            Inches(x0 / 72), Inches(y0 / 72), Inches((x1 - x0) / 72), Inches((y1 - y0) / 72))
                        tf = textbox.text_frame
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
                        p = tf.paragraphs[0]
                        p.text = span["text"]
                        p.font.size = Pt(span["size"])
                        p.font.name = self.default_font if self.enforce_default_font else span["font"]
                        p.font.color.rgb = RGBColor.from_string(f"{span['color']:06X}")
                        p.font.bold = bool(span["flags"] & 2**4)
                        p.font.italic = bool(span["flags"] & 2**1)
                        p.alignment = PP_ALIGN.LEFT

            # Add OCR text if scanned
            if scanned_document:
                page_ocr = ocr_results[page_num]
                for line in page_ocr[0]:
                    rect = line[0]
                    left = min(rect[0][0], rect[3][0]) * 0.24
                    top = min(rect[0][1], rect[1][1]) * 0.24
                    right = max(rect[1][0], rect[2][0]) * 0.24
                    bottom = max(rect[2][1], rect[3][1]) * 0.24
                    textbox = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(right-left), Pt(bottom-top))
                    tf = textbox.text_frame
                    tf.text = line[1][0]
                    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
                    for p in tf.paragraphs:
                        p.font.size = Pt(12)
                        if self.default_font:
                            p.font.name = self.default_font


# ------------------------------
# PDF -> PPTX service function
# ------------------------------

def convert_pdf_to_ppt(pdf_path: str, out_path: str, original_name: str,
                       enable_ocr: bool = False, default_font: str = "Arial"):
    try:
        # Set enable_ocr to False by default to disable OCR unless explicitly enabled
        converter = Converter(default_font=default_font, enable_ocr=enable_ocr, enforce_default_font=True)
        converter.convert(pdf_path, out_path)

        filename_suffix = "_ocr" if enable_ocr else ""
        # Encode the PPTX file into base64 for easy storage or transfer
        encoded_file = encode_file_to_base64(out_path)

        return {
            "success": True,
            "filename": Path(original_name).stem + filename_suffix + ".pptx",
            "file": encoded_file
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }
